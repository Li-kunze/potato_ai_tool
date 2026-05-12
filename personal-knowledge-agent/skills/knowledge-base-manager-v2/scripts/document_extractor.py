#!/usr/bin/env python3
"""
文档提取器 - 从复杂格式文件中提取结构化内容
供 knowledge-base-manager-agent 调用
"""

import os
import json
from pathlib import Path
from typing import Dict, Any, List, Optional


def detect_file_type(file_path: str) -> str:
    """检测文件类型"""
    ext = Path(file_path).suffix.lower()
    
    type_mapping = {
        '.docx': 'word',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
        '.pdf': 'pdf'
    }
    
    return type_mapping.get(ext, 'unknown')


def extract_word(file_path: str) -> Dict[str, Any]:
    """提取Word文档内容"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        
        # 提取段落
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 提取表格
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        
        # 纯文本内容（用于嵌入）
        content = '\n\n'.join(paragraphs)
        
        # 表格也转换为文本
        if tables:
            content += '\n\n[表格内容]\n'
            for i, table in enumerate(tables, 1):
                content += f'\n表{i}:\n'
                for row in table:
                    content += ' | '.join(row) + '\n'
        
        return {
            'success': True,
            'content': content,
            'metadata': {
                'title': doc.core_properties.title or Path(file_path).stem,
                'author': doc.core_properties.author or '',
                'type': 'word',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'paragraphs_count': len(paragraphs),
                'tables_count': len(tables),
                'word_count': len(content)
            },
            'structured_data': {
                'paragraphs': paragraphs,
                'tables': tables
            }
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'Word extraction failed: {str(e)}',
            'content': '',
            'metadata': {'type': 'word'}
        }


def extract_excel(file_path: str) -> Dict[str, Any]:
    """提取Excel文档内容"""
    try:
        import pandas as pd
        
        xl = pd.ExcelFile(file_path)
        sheets_data = {}
        all_text = []
        total_rows = 0
        
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            df = df.dropna(how='all').dropna(axis=1, how='all')
            
            sheets_data[sheet_name] = {
                'headers': df.columns.tolist(),
                'data': df.to_dict(orient='records'),
                'row_count': len(df)
            }
            
            # 转换为文本
            all_text.append(f'=== {sheet_name} ===')
            all_text.append(df.to_string())
            all_text.append('')
            
            total_rows += len(df)
        
        return {
            'success': True,
            'content': '\n'.join(all_text),
            'metadata': {
                'title': Path(file_path).stem,
                'type': 'excel',
                'sheets': list(sheets_data.keys()),
                'sheets_count': len(sheets_data),
                'total_rows': total_rows
            },
            'structured_data': {
                'sheets': sheets_data
            }
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'Excel extraction failed: {str(e)}',
            'content': '',
            'metadata': {'type': 'excel'}
        }


def extract_powerpoint(file_path: str) -> Dict[str, Any]:
    """提取PowerPoint文档内容"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides = []
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                slide_content = '\n'.join(slide_texts)
                slides.append({
                    'slide_number': i,
                    'text': slide_content
                })
                all_text.append(f'=== Slide {i} ===')
                all_text.append(slide_content)
                all_text.append('')
        
        return {
            'success': True,
            'content': '\n'.join(all_text),
            'metadata': {
                'title': Path(file_path).stem,
                'type': 'powerpoint',
                'slides_count': len(prs.slides),
                'extracted_slides': len(slides)
            },
            'structured_data': {
                'slides': slides
            }
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'PowerPoint extraction failed: {str(e)}',
            'content': '',
            'metadata': {'type': 'powerpoint'}
        }


def extract_pdf(file_path: str) -> Dict[str, Any]:
    """提取PDF文档内容"""
    try:
        import pdfplumber
        
        pages = []
        tables = []
        all_text = []
        
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    pages.append({
                        'page_number': i,
                        'text': text
                    })
                    all_text.append(f'=== Page {i} ===')
                    all_text.append(text)
                    all_text.append('')
                
                # 提取表格
                page_tables = page.extract_tables()
                for table in page_tables:
                    tables.append({
                        'page': i,
                        'data': table
                    })
        
        return {
            'success': True,
            'content': '\n'.join(all_text),
            'metadata': {
                'title': Path(file_path).stem,
                'type': 'pdf',
                'pages_count': len(pages),
                'tables_count': len(tables)
            },
            'structured_data': {
                'pages': pages,
                'tables': tables
            }
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'PDF extraction failed: {str(e)}',
            'content': '',
            'metadata': {'type': 'pdf'}
        }


def extract_document(file_path: str) -> Dict[str, Any]:
    """
    提取文档内容的主入口
    
    返回标准格式:
    {
        "success": true/false,
        "content": "纯文本内容",
        "metadata": {...},
        "structured_data": {...},
        "error": "错误信息（如有）"
    }
    """
    if not os.path.exists(file_path):
        return {
            'success': False,
            'error': f'File not found: {file_path}',
            'content': '',
            'metadata': {'source': file_path}
        }
    
    file_type = detect_file_type(file_path)
    
    processors = {
        'word': extract_word,
        'excel': extract_excel,
        'powerpoint': extract_powerpoint,
        'pdf': extract_pdf
    }
    
    processor = processors.get(file_type)
    if processor:
        result = processor(file_path)
        result['metadata']['source'] = file_path
        result['metadata']['file_name'] = os.path.basename(file_path)
        return result
    else:
        return {
            'success': False,
            'error': f'Unsupported file type: {file_type}',
            'content': '',
            'metadata': {
                'source': file_path,
                'detected_type': file_type
            }
        }


def batch_extract(folder_path: str, recursive: bool = False) -> List[Dict[str, Any]]:
    """批量提取文件夹中的文档"""
    results = []
    
    supported_exts = ['.docx', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf']
    
    if recursive:
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                if any(file.lower().endswith(ext) for ext in supported_exts):
                    file_path = os.path.join(root, file)
                    result = extract_document(file_path)
                    results.append(result)
    else:
        for file in os.listdir(folder_path):
            if any(file.lower().endswith(ext) for ext in supported_exts):
                file_path = os.path.join(folder_path, file)
                result = extract_document(file_path)
                results.append(result)
    
    return results


def main():
    """命令行入口"""
    import sys
    
    if len(sys.argv) < 2:
        print("Document Extractor for Knowledge Base")
        print("\nUsage:")
        print("  python document_extractor.py <file_path>")
        print("  python document_extractor.py batch <folder_path> [--recursive]")
        sys.exit(1)
    
    if sys.argv[1] == 'batch':
        if len(sys.argv) < 3:
            print("Usage: python document_extractor.py batch <folder_path> [--recursive]")
            sys.exit(1)
        
        folder_path = sys.argv[2]
        recursive = '--recursive' in sys.argv
        
        results = batch_extract(folder_path, recursive)
        print(json.dumps(results, ensure_ascii=False, indent=2))
    else:
        file_path = sys.argv[1]
        result = extract_document(file_path)
        print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
